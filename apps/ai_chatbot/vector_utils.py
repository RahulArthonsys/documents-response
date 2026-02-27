"""
vector_utils.py — ChromaDB operations, embedding generation, text extraction,
document indexing, and semantic search for the AI Chatbot.

Based on: EMBEDDING_CHROMA_FULL_DOCUMENTATION.md
"""
import logging
import json
import re
import os
from datetime import datetime
from typing import Dict, List, Any, Tuple

import numpy as np
from django.conf import settings
from django.utils import timezone

logger = logging.getLogger(__name__)

# ──────────────────────────────────────────────────────────
# 1. Gemini API Key — singleton
# ────────────────────────────────────────────────────────
GEMINI_EMBEDDING_MODEL = "models/gemini-embedding-001"
GEMINI_EMBEDDING_DIM = 3072

_gemini_api_key = getattr(settings, 'GEMINI_API_KEY', None) or os.environ.get('GEMINI_API_KEY', '')
client = None  # kept for backward-compat references

if _gemini_api_key:
    logger.info("Gemini API key loaded successfully")
else:
    logger.warning("GEMINI_API_KEY not set — embeddings will be disabled")


# Singleton LangChain embeddings object (uses v1 API, not v1beta)
_lc_embeddings = None


def get_lc_embeddings():
    """Return singleton langchain_google_genai.GoogleGenerativeAIEmbeddings."""
    global _lc_embeddings
    if _lc_embeddings is None:
        if not _gemini_api_key:
            raise ValueError("GEMINI_API_KEY not set")
        from langchain_google_genai import GoogleGenerativeAIEmbeddings
        _lc_embeddings = GoogleGenerativeAIEmbeddings(
            model=GEMINI_EMBEDDING_MODEL,
            google_api_key=_gemini_api_key,
        )
    return _lc_embeddings


# ──────────────────────────────────────────────────────────
# 2. ChromaDB Client — singleton
# ──────────────────────────────────────────────────────────
_chroma_client = None


def get_chroma_client():
    """Get or create a singleton ChromaDB PersistentClient."""
    global _chroma_client
    if _chroma_client is None:
        try:
            import chromadb
            chroma_path = os.path.join(settings.PROJECT_ROOT, 'chroma_store')
            os.makedirs(chroma_path, exist_ok=True)

            try:
                _chroma_client = chromadb.PersistentClient(path=chroma_path)
                logger.info(f"ChromaDB PersistentClient created at: {chroma_path}")
            except Exception as e1:
                logger.warning(f"PersistentClient failed ({e1}), trying EphemeralClient")
                _chroma_client = chromadb.EphemeralClient()
                logger.info("ChromaDB EphemeralClient created (in-memory)")
        except Exception as e:
            logger.error(f"Failed to create ChromaDB client: {e}")
            raise
    return _chroma_client


# ──────────────────────────────────────────────────────────
# 3. Custom ChromaDB Embedding Function
# ──────────────────────────────────────────────────────────
class GeminiEmbeddingFunction:
    """Custom embedding function for ChromaDB using Gemini gemini-embedding-001.
    Produces 3072-dimension vectors."""

    def __init__(self):
        self._model = GEMINI_EMBEDDING_MODEL
        self._dimension = GEMINI_EMBEDDING_DIM

    def name(self) -> str:
        return "gemini-embedding-001"

    def __call__(self, input: List[str]) -> List[List[float]]:
        """Called by ChromaDB when it needs embeddings for documents or queries."""
        return get_lc_embeddings().embed_documents(input)


_gemini_embedding_function = None


def get_openai_embedding_function():
    """Get or create the singleton Gemini embedding function (name kept for backward compat)."""
    global _gemini_embedding_function
    if _gemini_embedding_function is None:
        _gemini_embedding_function = GeminiEmbeddingFunction()
    return _gemini_embedding_function


# ──────────────────────────────────────────────────────────
# 4. Single & Batch Embedding Generation
# ──────────────────────────────────────────────────────────
def generate_gemini_embedding(text: str) -> List[float]:
    """Generate a single 3072-dim embedding vector using Gemini gemini-embedding-001."""
    return get_lc_embeddings().embed_query(text)


# Alias for backward compatibility
generate_openai_embedding = generate_gemini_embedding


def generate_gemini_embeddings_batch(texts: List[str], batch_size: int = 100) -> List[List[float]]:
    """Generate Gemini embeddings in batches."""
    lc_emb = get_lc_embeddings()
    all_embeddings = []
    for i in range(0, len(texts), batch_size):
        batch = texts[i:i + batch_size]
        try:
            batch_embeddings = lc_emb.embed_documents(batch)
            all_embeddings.extend(batch_embeddings)
            logger.info(f"Embedded batch {i // batch_size + 1} ({len(batch)} texts)")
        except Exception as e:
            logger.error(f"Batch embedding error at index {i}: {e}")
            all_embeddings.extend([[0.0] * GEMINI_EMBEDDING_DIM] * len(batch))
    return all_embeddings


# Alias for backward compatibility
generate_openai_embeddings_batch = generate_gemini_embeddings_batch


# ──────────────────────────────────────────────────────────
# 5. ChromaDB Collections
# ──────────────────────────────────────────────────────────
def _collection_has_wrong_dims(collection, expected_dim: int = GEMINI_EMBEDDING_DIM) -> bool:
    """Return True if the collection's stored vectors have different dimensions."""
    try:
        if collection.count() == 0:
            return False
        # Peek at one stored embedding
        peek = collection.get(limit=1, include=["embeddings"])
        embeddings = peek.get("embeddings")
        if embeddings and len(embeddings) > 0 and embeddings[0] is not None:
            stored_dim = len(embeddings[0])
            if stored_dim != expected_dim:
                logger.warning(
                    f"Collection '{collection.name}' has dim={stored_dim}, "
                    f"expected {expected_dim} — will recreate."
                )
                return True
    except Exception as e:
        logger.warning(f"Dim check failed for '{collection.name}': {e}")
    return False


def get_or_create_collection():
    """Get/create the 'documents' knowledge base collection with Gemini embeddings.
    Deletes and recreates if embedding dimensions conflict."""
    chroma_client = get_chroma_client()
    embedding_fn = get_openai_embedding_function()  # returns GeminiEmbeddingFunction

    def _make_collection():
        return chroma_client.get_or_create_collection(
            name="documents",
            embedding_function=embedding_fn,
            metadata={"hnsw:space": "cosine"}
        )

    try:
        collection = _make_collection()
        if _collection_has_wrong_dims(collection):
            chroma_client.delete_collection("documents")
            collection = _make_collection()
        return collection
    except Exception as conflict:
        logger.warning(f"Collection conflict for 'documents' — recreating: {conflict}")
        try:
            chroma_client.delete_collection("documents")
        except Exception:
            pass
        return _make_collection()


def get_session_collection_name(conversation_id: int) -> str:
    """Generate collection name for a session."""
    return f"session_conv_{conversation_id}"


def get_or_create_session_collection(conversation_id: int):
    """Get or create a session-specific ChromaDB collection with Gemini embeddings.
    Deletes and recreates on dimension conflict."""
    chroma_client = get_chroma_client()
    collection_name = get_session_collection_name(conversation_id)
    embedding_fn = get_openai_embedding_function()  # returns GeminiEmbeddingFunction

    def _make_collection():
        return chroma_client.get_or_create_collection(
            name=collection_name,
            embedding_function=embedding_fn,
            metadata={"hnsw:space": "cosine"}
        )

    try:
        collection = _make_collection()
        if _collection_has_wrong_dims(collection):
            chroma_client.delete_collection(collection_name)
            collection = _make_collection()
        return collection
    except Exception as conflict:
        logger.warning(f"Collection conflict for {collection_name} — recreating: {conflict}")
        try:
            chroma_client.delete_collection(collection_name)
        except Exception:
            pass
        return _make_collection()


# ──────────────────────────────────────────────────────────
# 6. Text Extraction (Multi-Format)
# ──────────────────────────────────────────────────────────
def extract_text_from_file(file_path: str, file_type: str) -> str:
    """Route to the correct extraction method based on file_type."""
    ft = file_type.lower().strip('.')

    if ft in ['pdf']:
        return extract_text(file_path)
    elif ft in ['txt', 'text', 'md']:
        with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
            return f.read()
    elif ft in ['docx']:
        try:
            import docx
            doc = docx.Document(file_path)
            return "\n".join([para.text for para in doc.paragraphs])
        except Exception as e:
            logger.error(f"DOCX extraction error: {e}")
            return ""
    elif ft in ['pptx', 'ppt']:
        try:
            from pptx import Presentation
            prs = Presentation(file_path)
            text_parts = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        text_parts.append(shape.text)
            return "\n".join(text_parts)
        except Exception as e:
            logger.error(f"PPTX extraction error: {e}")
            return ""
    elif ft in ['xlsx', 'xls']:
        try:
            import pandas as pd
            dfs = pd.read_excel(file_path, sheet_name=None)
            text_parts = []
            for sheet_name, df in dfs.items():
                text_parts.append(f"Sheet: {sheet_name}\n{df.to_string()}")
            return "\n\n".join(text_parts)
        except Exception as e:
            logger.error(f"Excel extraction error: {e}")
            return ""
    elif ft in ['csv']:
        try:
            import pandas as pd
            df = pd.read_csv(file_path)
            return df.to_string()
        except Exception as e:
            logger.error(f"CSV extraction error: {e}")
            return ""
    elif ft in ['json']:
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                data = json.load(f)
            return json.dumps(data, indent=2)
        except Exception as e:
            logger.error(f"JSON extraction error: {e}")
            return ""
    elif ft in ['xml']:
        try:
            import xml.etree.ElementTree as ET
            tree = ET.parse(file_path)
            root = tree.getroot()

            def extract_xml_text(element):
                texts = []
                if element.text and element.text.strip():
                    texts.append(element.text.strip())
                for child in element:
                    texts.extend(extract_xml_text(child))
                if element.tail and element.tail.strip():
                    texts.append(element.tail.strip())
                return texts

            return "\n".join(extract_xml_text(root))
        except Exception as e:
            logger.error(f"XML extraction error: {e}")
            return ""
    else:
        # Fallback: try reading as text
        try:
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                return f.read()
        except Exception:
            return ""


def extract_text(file_path: str) -> str:
    """Extract text from a PDF using PyMuPDF."""
    try:
        import fitz  # PyMuPDF
        doc = fitz.open(file_path)
        text_content = []
        for page in doc:
            text_content.append(page.get_text())
        doc.close()
        return "\n".join(text_content)
    except Exception as e:
        logger.error(f"PDF extraction error: {e}")
        return ""


def extract_text_and_images(file_path: str) -> dict:
    """Extract text, images, and tables from a PDF using PyMuPDF."""
    try:
        import fitz
        doc = fitz.open(file_path)
        text_content = []
        images = []
        tables = []

        for page_num in range(len(doc)):
            page = doc[page_num]
            text_content.append(page.get_text())

            # Extract images
            image_list = page.get_images(full=True)
            for img_idx, img_info in enumerate(image_list):
                try:
                    xref = img_info[0]
                    base_image = doc.extract_image(xref)
                    if base_image:
                        images.append({
                            "page": page_num + 1,
                            "image_data": base_image["image"],
                            "ext": base_image["ext"]
                        })
                except Exception:
                    pass

            # Try to extract tables
            try:
                page_tables = page.find_tables()
                if page_tables and page_tables.tables:
                    for table in page_tables.tables:
                        table_data = table.extract()
                        if table_data:
                            tables.append({
                                "page": page_num + 1,
                                "data": table_data
                            })
            except Exception:
                pass

        doc.close()
        return {
            "text": "\n".join(text_content),
            "images": images,
            "tables": tables
        }
    except Exception as e:
        logger.error(f"PDF extraction with images error: {e}")
        return {"text": "", "images": [], "tables": []}


# ──────────────────────────────────────────────────────────
# 7. Intelligent Chunking
# ──────────────────────────────────────────────────────────
def extract_metrics_from_text(text: str) -> List[Dict]:
    """Extract financial metrics and capacity scores from text using regex."""
    metrics = []

    financial_patterns = {
        "revenue": r"(?:revenue|sales)\s*(?:for|in)?\s*([Q\d\w\s]{1,15}?)[:=]\s*\$?(\d+(?:,\d{3})*(?:\.\d+)?)\s*(million|billion|m|b|k)?\b",
        "ebitda": r"(?:ebitda|earnings)\s*(?:for|in)?\s*([Q\d\w\s]{1,15}?)[:=]\s*\$?(\d+(?:,\d{3})*(?:\.\d+)?)\s*(million|billion|m|b|k)?\b",
        "profit": r"(?:profit|margin|income)\s*(?:for|in)?\s*([Q\d\w\s]{1,15}?)[:=]\s*\$?(\d+(?:,\d{3})*(?:\.\d+)?)\s*(million|billion|m|b|k)?(?:\s*%)?",
        "valuation": r"(?:valuation|value|worth)\s*(?:of|at)?\s*\$?(\d+(?:,\d{3})*(?:\.\d+)?)\s*(million|billion|m|b|k)?\b",
    }

    capacity_patterns = {
        "strategic_capacity": r"(?:strategic\s*capacity|capacity\s*score)\s*[:=]?\s*(\d+(?:\.\d+)?)",
        "overall_score": r"(?:overall\s*score|total\s*score)\s*[:=]?\s*(\d+(?:\.\d+)?)",
    }

    for metric_type, pattern in financial_patterns.items():
        for match in re.finditer(pattern, text, re.IGNORECASE):
            metrics.append({
                "type": metric_type,
                "value": match.group(0),
                "raw": match.group(0)
            })

    for metric_type, pattern in capacity_patterns.items():
        for match in re.finditer(pattern, text, re.IGNORECASE):
            metrics.append({
                "type": metric_type,
                "value": match.group(1),
                "raw": match.group(0)
            })

    return metrics


def intelligent_chunk_text(text: str, max_chunk_size: int = 800) -> List[Dict]:
    """Split text into chunks based on paragraph boundaries with metrics metadata."""
    paragraphs = re.split(r'\n\s*\n', text)
    chunks = []
    current_chunk = ""
    current_metrics = []

    for paragraph in paragraphs:
        paragraph = paragraph.strip()
        if not paragraph:
            continue

        para_metrics = extract_metrics_from_text(paragraph)

        if len(current_chunk) + len(paragraph) > max_chunk_size and current_chunk:
            chunks.append({
                "text": current_chunk.strip(),
                "metrics": current_metrics,
                "has_metrics": len(current_metrics) > 0
            })
            current_chunk = paragraph
            current_metrics = para_metrics
        else:
            current_chunk += "\n\n" + paragraph if current_chunk else paragraph
            current_metrics.extend(para_metrics)

    # Final chunk
    if current_chunk.strip():
        chunks.append({
            "text": current_chunk.strip(),
            "metrics": current_metrics,
            "has_metrics": len(current_metrics) > 0
        })

    return chunks


def format_table_as_text(table_data: list, page: int) -> str:
    """Convert a table (list of lists) to text."""
    if not table_data:
        return ""
    lines = []
    for row in table_data:
        cleaned = [str(cell) if cell else "" for cell in row]
        lines.append(" | ".join(cleaned))
    return f"[Table from page {page}]\n" + "\n".join(lines)


# ──────────────────────────────────────────────────────────
# 8. Knowledge Base Document Processing & Indexing
# ──────────────────────────────────────────────────────────
def process_document_content(document) -> dict:
    """Process a KnowledgeDocument: extract text, images, tables → chunks."""
    file_path = document.file.path
    file_ext = os.path.splitext(file_path)[1].lower().strip('.')

    all_chunks = []

    if file_ext == 'pdf':
        result = extract_text_and_images(file_path)
        text = result["text"]
        tables = result.get("tables", [])

        # Intelligent chunking of the main text
        text_chunks = intelligent_chunk_text(text, max_chunk_size=800)
        all_chunks.extend(text_chunks)

        # Table chunks
        for table in tables:
            table_text = format_table_as_text(table["data"], table["page"])
            if table_text.strip():
                all_chunks.append({
                    "text": table_text,
                    "metrics": extract_metrics_from_text(table_text),
                    "has_metrics": True,
                    "is_table": True,
                    "page": table["page"]
                })
    else:
        text = extract_text_from_file(file_path, file_ext)
        text_chunks = intelligent_chunk_text(text, max_chunk_size=800)
        all_chunks.extend(text_chunks)

    # Save extracted text to the document model
    extracted_text = text if 'text' in locals() else ""
    if extracted_text:
        document.content = extracted_text
        document.save(update_fields=['content'])

    return {
        "success": True if all_chunks else False,
        "chunks": all_chunks,
        "text": extracted_text,
        "chunk_count": len(all_chunks),
        "error": None if all_chunks else "No text could be extracted from the file"
    }


def index_document_embeddings(collection, document, chunks=None):
    """Index a KnowledgeDocument into the 'documents' ChromaDB collection."""
    try:
        # Step 1: Process document if no chunks provided
        if chunks is None:
            processed = process_document_content(document)
            chunk_data = processed["chunks"]
            chunks = [c["text"] for c in chunk_data]
        else:
            chunk_data = [{"text": c, "metrics": [], "has_metrics": False} for c in chunks]

        if not chunks:
            logger.warning(f"No chunks extracted from document {document.id}")
            return

        # Step 2: Build IDs and metadata
        ids = [f"{document.id}_{i}" for i in range(len(chunks))]
        metadatas = []
        for i, chunk_info in enumerate(chunk_data):
            metadata = {
                "document_id": str(document.id),
                "document_title": document.title,
                "chunk_index": i,
                "has_metrics": chunk_info.get("has_metrics", False),
                "is_table": chunk_info.get("is_table", False),
                "page": chunk_info.get("page", 0),
                "indexed_at": datetime.now().isoformat()
            }
            metadatas.append(metadata)

        # Step 3: Generate embeddings
        embeddings = generate_openai_embeddings_batch(chunks)

        # Step 4: Upsert into collection (batch if needed)
        batch_size = 100
        for i in range(0, len(chunks), batch_size):
            end = min(i + batch_size, len(chunks))
            collection.upsert(
                documents=chunks[i:end],
                ids=ids[i:end],
                metadatas=metadatas[i:end],
                embeddings=embeddings[i:end]
            )

        # Step 5: Update document metadata
        document.embedding_metadata = json.dumps({
            "chunk_count": len(chunks),
            "status": "indexed",
            "embedding_model": "gemini-embedding-001",
            "indexed_at": datetime.now().isoformat()
        })
        document.is_processed = True
        document.save(update_fields=['embedding_metadata', 'is_processed'])

        logger.info(f"Indexed {len(chunks)} chunks for document {document.id} ({document.title})")

    except Exception as e:
        logger.error(f"Error indexing document {document.id}: {e}")
        document.embedding_metadata = json.dumps({
            "status": "error",
            "error": str(e),
            "timestamp": datetime.now().isoformat()
        })
        document.save(update_fields=['embedding_metadata'])
        raise


# ──────────────────────────────────────────────────────────
# 9. Session Document Indexing
# ──────────────────────────────────────────────────────────
def index_session_document(file_path: str, file_type: str, conversation_id: int, original_filename: str) -> dict:
    """Index a user-uploaded document into a session-specific collection."""
    try:
        from langchain_text_splitters import RecursiveCharacterTextSplitter

        # Extract text
        text_content = extract_text_from_file(file_path, file_type)
        if not text_content or not text_content.strip():
            return {"success": False, "error": "No text could be extracted from the file"}

        # Chunk with LangChain splitter
        text_splitter = RecursiveCharacterTextSplitter(
            chunk_size=1000,
            chunk_overlap=200,
            length_function=len
        )
        text_chunks = text_splitter.split_text(text_content)

        if not text_chunks:
            return {"success": False, "error": "No chunks generated"}

        # Get or create session collection
        collection = get_or_create_session_collection(conversation_id)

        # Build IDs and metadata
        safe_filename = re.sub(r'[^a-zA-Z0-9._-]', '_', original_filename)
        ids = [f"session_{conversation_id}_{safe_filename}_{i}" for i in range(len(text_chunks))]
        metadatas = [{
            "source": original_filename,
            "conversation_id": str(conversation_id),
            "chunk_index": i,
            "file_type": file_type,
            "is_table": False,
            "page": 0,
            "upload_date": datetime.now().isoformat()
        } for i in range(len(text_chunks))]

        # Generate embeddings
        embeddings = generate_openai_embeddings_batch(text_chunks)

        # Upsert into collection
        batch_size = 100
        for i in range(0, len(text_chunks), batch_size):
            end = min(i + batch_size, len(text_chunks))
            collection.upsert(
                documents=text_chunks[i:end],
                ids=ids[i:end],
                metadatas=metadatas[i:end],
                embeddings=embeddings[i:end]
            )

        logger.info(f"Indexed {len(text_chunks)} session chunks for conversation {conversation_id}")
        return {
            "success": True,
            "chunks_created": len(text_chunks),
            "collection_name": get_session_collection_name(conversation_id)
        }

    except Exception as e:
        logger.error(f"Session indexing error: {e}")
        return {"success": False, "error": str(e)}


# ──────────────────────────────────────────────────────────
# 10. Search & Retrieval
# ──────────────────────────────────────────────────────────
def search_documents(query: str, collection, top_k: int = 5, metadata_filter: dict = None) -> dict:
    """Semantic search on the 'documents' (KB) collection using OpenAI embeddings."""
    try:
        query_embedding = generate_openai_embedding(query)

        query_params = {
            "query_embeddings": [query_embedding],
            "n_results": min(top_k, collection.count()) if collection.count() > 0 else top_k,
            "include": ["documents", "metadatas", "distances"]
        }
        if metadata_filter:
            query_params["where"] = metadata_filter

        results = collection.query(**query_params)

        search_results = {
            "documents": results.get("documents", [[]])[0],
            "metadatas": results.get("metadatas", [[]])[0],
            "distances": results.get("distances", [[]])[0],
            "ids": results.get("ids", [[]])[0],
        }

        # Filter by distance threshold
        filtered_docs = []
        filtered_meta = []
        filtered_dist = []
        for doc, meta, dist in zip(
            search_results["documents"],
            search_results["metadatas"],
            search_results["distances"]
        ):
            if dist <= 1.5:  # Distance threshold
                filtered_docs.append(doc)
                filtered_meta.append(meta)
                filtered_dist.append(dist)

        return {
            "documents": filtered_docs,
            "metadatas": filtered_meta,
            "distances": filtered_dist,
        }

    except Exception as e:
        logger.error(f"Search error: {e}")
        return {"documents": [], "metadatas": [], "distances": []}


def search_session_documents(query: str, conversation_id: int, top_k: int = 15, specific_filename: str = None) -> list:
    """Search within a session-specific ChromaDB collection."""
    try:
        collection = get_or_create_session_collection(conversation_id)

        if collection.count() == 0:
            return []

        query_embedding = generate_openai_embedding(query)

        query_params = {
            "query_embeddings": [query_embedding],
            "n_results": min(top_k, collection.count()),
            "include": ["documents", "metadatas", "distances"]
        }

        if specific_filename:
            query_params["where"] = {"source": specific_filename}

        results = collection.query(**query_params)

        formatted_results = []
        docs = results.get("documents", [[]])[0]
        metas = results.get("metadatas", [[]])[0]
        dists = results.get("distances", [[]])[0]

        for doc, meta, dist in zip(docs, metas, dists):
            if dist <= 1.5:
                formatted_results.append({
                    "content": doc,
                    "metadata": meta,
                    "distance": dist,
                    "source": meta.get("source", "unknown")
                })

        return formatted_results

    except Exception as e:
        logger.error(f"Session search error: {e}")
        return []


# ──────────────────────────────────────────────────────────
# 11. Cleanup & Deletion
# ──────────────────────────────────────────────────────────
def delete_session_collection(conversation_id: int) -> bool:
    """Delete a session-specific ChromaDB collection."""
    try:
        chroma_client = get_chroma_client()
        collection_name = get_session_collection_name(conversation_id)
        chroma_client.delete_collection(collection_name)
        logger.info(f"Deleted session collection: {collection_name}")
        return True
    except Exception as e:
        logger.warning(f"Could not delete session collection {conversation_id}: {e}")
        return False


def delete_document_embeddings(collection, document_id: int, max_chunks: int = 1000):
    """Delete all embeddings for a KnowledgeDocument from ChromaDB."""
    try:
        ids_to_delete = [f"{document_id}_{i}" for i in range(max_chunks)]
        # ChromaDB silently ignores non-existent IDs
        collection.delete(ids=ids_to_delete)
        logger.info(f"Deleted embeddings for document {document_id}")
    except Exception as e:
        logger.error(f"Error deleting embeddings for document {document_id}: {e}")
